"""
org_chart_tool.py — Excel-to-PowerPoint organization chart generator

Reads a workbook where each sheet/tab is a department, with employee rows
containing Name / Designation / Department / ReportsToName columns, and
produces a fully laid-out PowerPoint organogram: one connected tree per
department, paginated automatically when a manager has too many direct
reports to fit legibly on one slide.

USAGE
-----
    python org_chart_tool.py INPUT.xlsx [-o OUTPUT.pptx] [--company "Name"]

    python org_chart_tool.py "Kaiser Permanente_Org_Charting.xlsx" \\
        -o "Kaiser_Permanente_Org_Chart.pptx" --company "Kaiser Permanente"

REQUIRED EXCEL COLUMNS (per sheet)
----------------------------------
    Name            - employee full name
    Designation     - job title
    Department      - sub-department / region label shown next to the title
    ReportsToName   - the Name of this person's manager, WITHIN THE SAME
                      SHEET. Blank / not found in the sheet => top-level
                      box for that department.

Any other columns in the workbook are ignored. Nothing in the source
workbook is read via formulas (values only), and the workbook is never
modified.

INSTALL
-------
    pip install openpyxl python-pptx
"""
import argparse
import datetime
import os
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

LEVEL_COLORS = ["0C3649", "347490", "4396BB", "70AFCC"]
LINE_GRAY = "9AA5AB"
NAVY = "0C3649"
GRAY_TEXT = "555555"
LINK_BLUE = "1F77C4"      # name text color, matching the hyperlink look in the reference
BODY_TEXT = "222222"      # title/department text under the name
CARD_FILL = "F5F8FB"      # soft light fill for each card
CONNECTOR_BLUE = "6FA8C9" # connector line color
SHADOW_TINT = "347490"    # backing "stacked card" shape behind each box
SHADOW_OFFSET = 0.13      # inches, how far the shadow card peeks out
CARD_CUT = 0.18           # inches, depth of the top-right/bottom-left corner cuts

SLIDE_W, SLIDE_H = 13.333, 7.5
LEFT_MARGIN, RIGHT_MARGIN = 0.45, 0.45
TOP, BOTTOM = 1.25, 0.4
USABLE_W = SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN
USABLE_H = SLIDE_H - TOP - BOTTOM

NODE_W, NODE_H = 1.7, 1.0
HGAP, VGAP = 0.32, 0.62
BLOCK_GAP = 0.5   # gap between independent blocks placed side by side on a page

FONT = "Arial"


class Emp:
    __slots__ = ("name", "title", "dept", "reports_to", "children", "level",
                 "x", "y", "subtree_w", "subtree_h", "linkedin")

    def __init__(self, name, title, dept, reports_to, linkedin=None):
        self.name = name
        self.title = title
        self.dept = dept
        self.reports_to = reports_to
        self.children = []
        self.level = 1
        self.x = 0.0
        self.y = 0.0
        self.subtree_w = NODE_W
        self.subtree_h = NODE_H
        self.linkedin = linkedin


REQUIRED_COLS = ("Name", "ReportsToName")
OPTIONAL_COLS = ("Designation", "Department", "LinkedinUrl")


def _clean_url(value):
    """Returns a usable URL string, or None for blanks/placeholders like 'NA'."""
    if not value:
        return None
    v = str(value).strip()
    if not v or v.upper() in ("NA", "N/A", "-", "NONE"):
        return None
    if not v.lower().startswith(("http://", "https://")):
        v = "https://" + v
    return v


def read_department(ws):
    headers = [c.value for c in ws[1]]
    col = {h: i for i, h in enumerate(headers) if h}
    missing = [c for c in REQUIRED_COLS if c not in col]
    if missing:
        raise ValueError(
            f"Sheet '{ws.title}' is missing required column(s): {', '.join(missing)}. "
            f"Found columns: {', '.join(str(h) for h in headers if h)}"
        )

    emps = []
    for row in ws.iter_rows(min_row=2, values_only=True):
        name = row[col["Name"]]
        if not name or not str(name).strip():
            continue
        title = (row[col["Designation"]] if "Designation" in col else "") or ""
        dept = (row[col["Department"]] if "Department" in col else "") or ""
        reports_to = row[col["ReportsToName"]]
        reports_to = str(reports_to).strip() if reports_to and str(reports_to).strip() else None
        linkedin = _clean_url(row[col["LinkedinUrl"]]) if "LinkedinUrl" in col else None
        emps.append(Emp(str(name).strip(), str(title).strip(), str(dept).strip(), reports_to, linkedin))

    by_name = {}
    for e in emps:
        by_name.setdefault(e.name, e)

    roots = []
    for e in emps:
        parent = by_name.get(e.reports_to) if e.reports_to else None
        if parent is not None and parent is not e:
            parent.children.append(e)
        else:
            roots.append(e)

    def set_level(e, lvl):
        e.level = lvl
        for c in e.children:
            set_level(c, lvl + 1)

    for r in roots:
        set_level(r, 1)

    return emps, roots


# ------------------------------------------------------------------
# Standard tidy-tree measure/layout (single row of children per node)
# ------------------------------------------------------------------
def measure(e):
    if not e.children:
        e.subtree_w = NODE_W
        e.subtree_h = NODE_H
        return e.subtree_w
    w = sum(measure(c) for c in e.children) + (len(e.children) - 1) * HGAP
    e.subtree_w = max(w, NODE_W)
    e.subtree_h = NODE_H + VGAP + max(c.subtree_h for c in e.children)
    return e.subtree_w


def layout(e, start_x, y):
    e.y = y
    if not e.children:
        e.x = start_x
        return
    cx = start_x
    first_slot = start_x
    last_slot = start_x
    for c in e.children:
        last_slot = cx
        layout(c, cx, y + NODE_H + VGAP)
        cx += c.subtree_w + HGAP
    # Center over the full bounding box of children's subtrees (their
    # allotted slots), NOT over children's own box positions -- a non-leaf
    # child's box is centered over *its* children and is generally not at
    # the left edge of its own subtree, so using child.x here would drift.
    left = first_slot
    right = last_slot + e.children[-1].subtree_w
    e.x = (left + right) / 2 - NODE_W / 2


def flatten(e, out):
    out.append(e)
    for c in e.children:
        flatten(c, out)


def pack_into_groups(items, max_width, gap):
    groups, current, current_w = [], [], 0.0
    for it in items:
        w = it.subtree_w
        add = w if not current else gap + w
        if current and current_w + add > max_width:
            groups.append(current)
            current, current_w = [it], w
        else:
            current.append(it)
            current_w += add
    if current:
        groups.append(current)
    return groups


# ------------------------------------------------------------------
# Build "blocks": a block is (node, subset_of_node's_direct_children)
# that fits within one slide's width. Any node (at any depth) whose
# children collectively overflow gets split into several blocks, each
# repeating that node; only the first group stays nested under its
# real parent, the rest bubble up as independent continuation blocks.
# ------------------------------------------------------------------
def prepare(node, budget):
    extra = []
    for c in node.children:
        extra += prepare(c, budget)

    if not node.children:
        node.subtree_w = NODE_W
        node.subtree_h = NODE_H
        return extra

    total_w = sum(c.subtree_w for c in node.children) + HGAP * (len(node.children) - 1)
    if total_w <= budget:
        node.subtree_w = max(total_w, NODE_W)
        node.subtree_h = NODE_H + VGAP + max(c.subtree_h for c in node.children)
        return extra

    groups = pack_into_groups(node.children, budget - NODE_W, HGAP)
    node.children = groups[0]
    node.subtree_w = block_width(node.children)
    node.subtree_h = NODE_H + VGAP + max(c.subtree_h for c in node.children)
    for g in groups[1:]:
        extra.append((node, g))
    return extra


def build_blocks(root):
    extra = prepare(root, USABLE_W)
    return [(root, list(root.children))] + extra


def block_width(children_subset):
    if not children_subset:
        return NODE_W
    return sum(c.subtree_w for c in children_subset) + HGAP * (len(children_subset) - 1)


def block_height(root, children_subset):
    if not children_subset:
        return NODE_H
    return NODE_H + VGAP + max(c.subtree_h for c in children_subset)


def layout_block(root, children_subset, x, root_y=0.0):
    """Position children_subset starting at left edge x (root may be a shared
    object repeated across several blocks, so its own position is *returned*
    rather than stored on the object, which would get clobbered by later
    blocks reusing the same root)."""
    cx = x
    first_slot = x
    last_slot = x
    for c in children_subset:
        last_slot = cx
        layout(c, cx, root_y + NODE_H + VGAP)
        cx += c.subtree_w + HGAP
    if children_subset:
        left = first_slot
        right = last_slot + children_subset[-1].subtree_w
        root_x = (left + right) / 2 - NODE_W / 2
    else:
        root_x = x
    return root_x


# ------------------------------------------------------------------
# Rendering
# ------------------------------------------------------------------
def set_run(run, text, size, bold, color_hex, font=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = RGBColor.from_string(color_hex)


def _card_points(x, y, w, h, cut):
    """Six points outlining a rectangle with its top-right and bottom-left
    corners cut off diagonally."""
    return [
        (x, y),
        (x + w - cut, y),
        (x + w, y + cut),
        (x + w, y + h),
        (x + cut, y + h),
        (x, y + h - cut),
    ]


def add_box(slide, x, y, w, h, e, scale=1.0):
    border_color = LEVEL_COLORS[min(e.level - 1, len(LEVEL_COLORS) - 1)]
    cut = CARD_CUT * scale

    # Backing "shadow" card: a slightly offset, darker copy of the same
    # cut-corner shape behind the main card, peeking out on the top-left,
    # giving the stacked/3D look from the reference screenshot.
    off = SHADOW_OFFSET * scale
    shadow = add_polygon(slide, _card_points(x - off, y - off, w, h, cut), fill_hex=SHADOW_TINT)

    shp = add_polygon(slide, _card_points(x, y, w, h, cut), fill_hex=CARD_FILL, line_hex=border_color)
    shp.line.width = Pt(1.0)

    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    name_size = max(8.5, 11.5 * scale)
    title_size = max(7.5, 9.5 * scale)

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    # Names are always styled like the reference's hyperlink look (blue,
    # underlined); only names with a real LinkedIn URL get an actual link.
    set_run(r1, e.name, name_size, True, LINK_BLUE)
    r1.font.underline = True
    if e.linkedin:
        r1.hyperlink.address = e.linkedin

    subtitle = e.title
    if e.dept and e.dept.lower() not in (e.title or "").lower():
        subtitle = f"{e.title}, {e.dept}" if e.title else e.dept
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        set_run(r2, subtitle, title_size, False, BODY_TEXT)


def add_line(slide, x1, y1, x2, y2, color=CONNECTOR_BLUE, weight=1.0):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                       Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = RGBColor.from_string(color)
    conn.line.width = Pt(weight)


def draw_tree_at(slide, e, ox, oy, scale):
    """Draw e and all its descendants using precomputed e.x/e.y, transformed by (ox, oy, scale)."""
    nodes = []
    flatten(e, nodes)
    for node in nodes:
        add_box(slide, ox + node.x * scale, oy + node.y * scale, NODE_W * scale, NODE_H * scale, node, scale)
    _draw_connectors_at(slide, e, ox, oy, scale)


def _draw_connectors_at(slide, e, ox, oy, scale):
    if not e.children:
        return
    px = ox + (e.x + NODE_W / 2) * scale
    pbottom = oy + (e.y + NODE_H) * scale
    bus_y = oy + (e.y + NODE_H + VGAP / 2) * scale
    add_line(slide, px, pbottom, px, bus_y)
    child_centers = [ox + (c.x + NODE_W / 2) * scale for c in e.children]
    # The bus line must span the parent's own drop point too, not just
    # first-to-last child: a child's box isn't necessarily centered under
    # its own subtree, so the parent's x and a lone/off-center child's x
    # can differ even when there's only one child.
    left_x = min(child_centers + [px])
    right_x = max(child_centers + [px])
    if left_x != right_x:
        add_line(slide, left_x, bus_y, right_x, bus_y)
    for c, ccx in zip(e.children, child_centers):
        ctop = oy + c.y * scale
        add_line(slide, ccx, bus_y, ccx, ctop)
        _draw_connectors_at(slide, c, ox, oy, scale)


def add_continuation_tag(slide, x, y, w, manager_name, scale):
    """Small caption + dashed stub shown above a box when that box's root
    is a continuation of someone who reports to a manager not shown on
    this slide (e.g. their own subtree was too large to fit alongside
    their manager, so they were promoted to their own block)."""
    tag_h = 0.22
    tag_y = max(0.05, y - tag_h - 0.03)

    stub = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                       Inches(x + w / 2), Inches(tag_y + tag_h),
                                       Inches(x + w / 2), Inches(y))
    stub.line.color.rgb = RGBColor.from_string(LINE_GRAY)
    stub.line.width = Pt(1.0)
    stub.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    tb = slide.shapes.add_textbox(Inches(x - 0.6), Inches(tag_y), Inches(w + 1.2), Inches(tag_h))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, f"continued \u2014 reports to {manager_name}", max(6.5, 7.5 * scale), False, GRAY_TEXT)
    for run in p.runs:
        run.font.italic = True


def draw_block_at(slide, root, children_subset, root_x, root_y, ox, oy, scale):
    """Draw a block using positions already computed by layout_block, transformed by (ox, oy, scale)."""
    box_x = ox + root_x * scale
    box_y = oy + root_y * scale
    add_box(slide, box_x, box_y, NODE_W * scale, NODE_H * scale, root, scale)
    if root.reports_to:
        add_continuation_tag(slide, box_x, box_y, NODE_W * scale, root.reports_to, scale)
    if children_subset:
        px = ox + (root_x + NODE_W / 2) * scale
        pbottom = oy + (root_y + NODE_H) * scale
        bus_y = oy + (root_y + NODE_H + VGAP / 2) * scale
        add_line(slide, px, pbottom, px, bus_y)
        child_centers = [ox + (c.x + NODE_W / 2) * scale for c in children_subset]
        left_x = min(child_centers + [px])
        right_x = max(child_centers + [px])
        if left_x != right_x:
            add_line(slide, left_x, bus_y, right_x, bus_y)
        for c, ccx in zip(children_subset, child_centers):
            ctop = oy + c.y * scale
            add_line(slide, ccx, bus_y, ccx, ctop)
            draw_tree_at(slide, c, ox, oy, scale)


_FREEFORM_UNIT = Inches(1) / 1000  # 1000 local units per inch, for sub-integer precision


def add_polygon(slide, points_in, fill_hex=None, line_hex=None):
    """Draw a closed polygon from a list of (x, y) points given in inches."""
    pts = [(round(x * 1000), round(y * 1000)) for x, y in points_in]
    fb = slide.shapes.build_freeform(pts[0][0], pts[0][1], scale=_FREEFORM_UNIT)
    fb.add_line_segments(pts[1:], close=True)
    shape = fb.convert_to_shape(0, 0)
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    else:
        shape.fill.background()
    if line_hex:
        shape.line.color.rgb = RGBColor.from_string(line_hex)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_title(slide, text):
    """Plain slide title text, no banner background."""
    tb = slide.shapes.add_textbox(Inches(LEFT_MARGIN), Inches(0.25),
                                   Inches(SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN), Inches(0.6))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    r = p.add_run()
    set_run(r, text, 20, True, "FFFFFF")


def render_department(prs, dept_title, roots):
    # gather all blocks across all roots of this department
    blocks = []
    for root in roots:
        blocks.extend(build_blocks(root))

    # pack blocks into pages (by width)
    pages, current, current_w = [], [], 0.0
    for root, subset in blocks:
        w = block_width(subset)
        add = w if not current else BLOCK_GAP + w
        if current and current_w + add > USABLE_W:
            pages.append(current)
            current, current_w = [(root, subset)], w
        else:
            current.append((root, subset))
            current_w += add
    if current:
        pages.append(current)

    total_pages = len(pages)
    for idx, page_blocks in enumerate(pages, start=1):
        label = f"{dept_title} ({idx}/{total_pages})" if total_pages > 1 else dept_title
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(slide, label)

        widths = [block_width(subset) for _, subset in page_blocks]
        heights = [block_height(root, subset) for root, subset in page_blocks]
        total_w = sum(widths) + BLOCK_GAP * (len(page_blocks) - 1)
        max_h = max(heights) if heights else NODE_H

        # First pass: compute every node's position at a provisional origin,
        # without drawing anything yet. Each block's root position is kept
        # in a local list (root objects can repeat across blocks when a
        # manager has too many direct reports, so we must not store the
        # position back on the shared object).
        cursor = 0.0
        root_positions = []
        all_bounds = []  # (x, y, w, h) rectangles, in provisional coordinates
        for (root, subset), w in zip(page_blocks, widths):
            root_x = layout_block(root, subset, cursor, root_y=0.0)
            root_positions.append((root_x, 0.0))
            all_bounds.append((root_x, 0.0, NODE_W, NODE_H))
            for c in subset:
                nodes = []
                flatten(c, nodes)
                for node in nodes:
                    all_bounds.append((node.x, node.y, NODE_W, NODE_H))
            cursor += w + BLOCK_GAP

        min_x = min(bx for bx, by, bw, bh in all_bounds)
        max_x = max(bx + bw for bx, by, bw, bh in all_bounds)
        min_y = min(by for bx, by, bw, bh in all_bounds)
        max_y = max(by + bh for bx, by, bw, bh in all_bounds)
        actual_w = max_x - min_x
        actual_h = max_y - min_y

        # Pages with a continuation tag (a promoted node repeated without
        # its real manager) need a little extra headroom for the dashed
        # stub + "reports to" caption above the top row of boxes.
        needs_tag_room = any(root.reports_to for root, _ in page_blocks)
        page_top = TOP + (0.3 if needs_tag_room else 0.0)
        page_usable_h = SLIDE_H - page_top - BOTTOM

        fit_scale = min(1.0, USABLE_W / actual_w if actual_w > 0 else 1.0,
                         page_usable_h / actual_h if actual_h > 0 else 1.0)

        ox = LEFT_MARGIN + max(0, (USABLE_W - actual_w * fit_scale) / 2) - min_x * fit_scale
        oy = page_top + max(0, (page_usable_h - actual_h * fit_scale) / 2) - min_y * fit_scale

        for (root, subset), (root_x, root_y) in zip(page_blocks, root_positions):
            draw_block_at(slide, root, subset, root_x, root_y, ox, oy, fit_scale)


def add_title_slide(prs, company_name, subtitle="Organizational Chart"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(NAVY)
    bg.line.fill.background()
    bg.shadow.inherit = False

    tb = slide.shapes.add_textbox(Inches(0), Inches(2.9), Inches(SLIDE_W), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, company_name, 40, True, "FFFFFF")

    tb2 = slide.shapes.add_textbox(Inches(0), Inches(3.75), Inches(SLIDE_W), Inches(0.7))
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    set_run(r2, subtitle, 22, False, "70AFCC")

    tb3 = slide.shapes.add_textbox(Inches(0), Inches(4.35), Inches(SLIDE_W), Inches(0.5))
    p3 = tb3.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    set_run(r3, datetime.date.today().strftime("Prepared %B %Y"), 13, False, "9AA5AB")


DISCLAIMER_LINES = [
    "Organograms are a graphical representation of an organization\u2019s hierarchy. Their purpose is to illustrate "
    "the reporting relationships and chain of command within an organization.",
    "This deck was generated directly from the source workbook: one department per tab, using each tab\u2019s "
    "\u201cReportsToName\u201d column to build the reporting tree, with no manual changes to the underlying data.",
    "Each person\u2019s title line shows their Designation and Department exactly as recorded in the workbook.",
    "Where a person\u2019s manager is not listed within the same department tab, that person is shown as a "
    "top-level box for that department.",
    "Departments with a large number of direct reports are split across multiple slides; the manager is repeated "
    "at the top of each continuation slide.",
]


def add_disclaimer_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Disclaimer")
    tb = slide.shapes.add_textbox(Inches(LEFT_MARGIN), Inches(1.5),
                                   Inches(SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN), Inches(3.4))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(DISCLAIMER_LINES):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        r = p.add_run()
        set_run(r, "\u2022  " + line, 13, False, "333333")


def build(input_path, output_path, company_name):
    wb = openpyxl.load_workbook(input_path, data_only=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    add_title_slide(prs, company_name)
    add_disclaimer_slide(prs)

    total_boxes = 0
    for ws in wb.worksheets:
        emps, roots = read_department(ws)
        if not emps:
            continue
        render_department(prs, ws.title, roots)
        total_boxes += len(emps)
        print(f"  {ws.title}: {len(emps)} employees, {len(roots)} top-level box(es)")

    prs.save(output_path)
    print(f"\nTotal employees plotted : {total_boxes}")
    print(f"Total slides            : {len(prs.slides._sldIdLst)}")
    print(f"Saved to                : {output_path}")


def main():
    parser = argparse.ArgumentParser(
        description="Generate a PowerPoint organization chart from an Excel workbook "
                    "(one department per sheet).",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    parser.add_argument("input", help="Path to the source .xlsx workbook")
    parser.add_argument("-o", "--output", default=None,
                         help="Path for the generated .pptx (default: <input name>_Org_Chart.pptx)")
    parser.add_argument("--company", default=None,
                         help="Company name shown on the title slide (default: derived from the input filename)")
    args = parser.parse_args()

    if not os.path.isfile(args.input):
        raise SystemExit(f"Input file not found: {args.input}")

    output_path = args.output or (os.path.splitext(os.path.basename(args.input))[0] + "_Org_Chart.pptx")
    company_name = args.company or os.path.splitext(os.path.basename(args.input))[0].replace("_", " ")

    print(f"Reading workbook: {args.input}")
    build(args.input, output_path, company_name)


if __name__ == "__main__":
    main()
